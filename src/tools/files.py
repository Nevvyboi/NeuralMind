"""
GroundZero AI - File Creation Engine
=====================================

Create professional documents:
1. Word documents (.docx)
2. PDF files
3. Excel spreadsheets (.xlsx)
4. PowerPoint presentations (.pptx)
5. CSV files
6. Markdown files
7. HTML files
"""

import os
import json
import csv
import io
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Union
from dataclasses import dataclass, field

try:
    from ..utils import get_data_path, ensure_dir, logger, timestamp, generate_id
except ImportError:
    from utils import get_data_path, ensure_dir, logger, timestamp, generate_id


# ============================================================================
# WORD DOCUMENT CREATOR
# ============================================================================

class WordCreator:
    """Create Word documents (.docx)."""
    
    @staticmethod
    def create(
        filepath: str,
        content: Union[str, List[Dict]],
        title: str = None,
        author: str = "GroundZero AI",
    ) -> str:
        """
        Create a Word document.
        
        Args:
            filepath: Output path
            content: Either a string or list of elements:
                     [{"type": "heading", "text": "...", "level": 1},
                      {"type": "paragraph", "text": "..."},
                      {"type": "table", "headers": [...], "rows": [...]},
                      {"type": "list", "items": [...]}]
            title: Document title
            author: Document author
        
        Returns:
            Path to created file
        """
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")
        
        doc = Document()
        
        # Set metadata
        doc.core_properties.author = author
        if title:
            doc.core_properties.title = title
            # Add title
            title_para = doc.add_heading(title, 0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Process content
        if isinstance(content, str):
            # Simple string content - split by double newlines for paragraphs
            for para in content.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())
        else:
            # Structured content
            for element in content:
                elem_type = element.get("type", "paragraph")
                
                if elem_type == "heading":
                    level = element.get("level", 1)
                    doc.add_heading(element.get("text", ""), level)
                
                elif elem_type == "paragraph":
                    para = doc.add_paragraph(element.get("text", ""))
                    if element.get("bold"):
                        for run in para.runs:
                            run.bold = True
                    if element.get("italic"):
                        for run in para.runs:
                            run.italic = True
                
                elif elem_type == "table":
                    headers = element.get("headers", [])
                    rows = element.get("rows", [])
                    
                    if headers or rows:
                        table = doc.add_table(rows=1 + len(rows), cols=len(headers or rows[0]))
                        table.style = 'Table Grid'
                        
                        # Add headers
                        if headers:
                            for i, header in enumerate(headers):
                                table.rows[0].cells[i].text = str(header)
                        
                        # Add rows
                        for row_idx, row in enumerate(rows):
                            for col_idx, cell in enumerate(row):
                                table.rows[row_idx + 1].cells[col_idx].text = str(cell)
                
                elif elem_type == "list":
                    for item in element.get("items", []):
                        doc.add_paragraph(str(item), style='List Bullet')
                
                elif elem_type == "numbered_list":
                    for item in element.get("items", []):
                        doc.add_paragraph(str(item), style='List Number')
                
                elif elem_type == "page_break":
                    doc.add_page_break()
        
        # Save
        ensure_dir(Path(filepath).parent)
        doc.save(filepath)
        
        logger.info(f"Created Word document: {filepath}")
        return filepath


# ============================================================================
# PDF CREATOR
# ============================================================================

class PDFCreator:
    """Create PDF files."""
    
    @staticmethod
    def create(
        filepath: str,
        content: Union[str, List[Dict]],
        title: str = None,
        author: str = "GroundZero AI",
    ) -> str:
        """
        Create a PDF document.
        
        Args:
            filepath: Output path
            content: String or list of elements
            title: Document title
            author: Document author
        
        Returns:
            Path to created file
        """
        # Try reportlab first
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
            from reportlab.lib.enums import TA_CENTER
            
            doc = SimpleDocTemplate(
                filepath,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            # Add title
            if title:
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Heading1'],
                    alignment=TA_CENTER,
                    spaceAfter=30,
                )
                story.append(Paragraph(title, title_style))
            
            # Process content
            if isinstance(content, str):
                for para in content.split('\n\n'):
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                        story.append(Spacer(1, 12))
            else:
                for element in content:
                    elem_type = element.get("type", "paragraph")
                    
                    if elem_type == "heading":
                        level = element.get("level", 1)
                        style = styles[f'Heading{min(level, 3)}']
                        story.append(Paragraph(element.get("text", ""), style))
                        story.append(Spacer(1, 12))
                    
                    elif elem_type == "paragraph":
                        story.append(Paragraph(element.get("text", ""), styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    elif elem_type == "table":
                        headers = element.get("headers", [])
                        rows = element.get("rows", [])
                        
                        data = []
                        if headers:
                            data.append(headers)
                        data.extend(rows)
                        
                        if data:
                            t = Table(data)
                            t.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                            ]))
                            story.append(t)
                            story.append(Spacer(1, 12))
                    
                    elif elem_type == "list":
                        for item in element.get("items", []):
                            story.append(Paragraph(f"• {item}", styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    elif elem_type == "page_break":
                        story.append(PageBreak())
            
            doc.build(story)
            logger.info(f"Created PDF: {filepath}")
            return filepath
            
        except ImportError:
            pass
        
        # Try fpdf as fallback
        try:
            from fpdf import FPDF
            
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            if title:
                pdf.set_font("Arial", 'B', 16)
                pdf.cell(0, 10, title, ln=True, align='C')
                pdf.set_font("Arial", size=12)
                pdf.ln(10)
            
            if isinstance(content, str):
                pdf.multi_cell(0, 10, content)
            else:
                for element in content:
                    elem_type = element.get("type", "paragraph")
                    text = element.get("text", "")
                    
                    if elem_type == "heading":
                        pdf.set_font("Arial", 'B', 14)
                        pdf.cell(0, 10, text, ln=True)
                        pdf.set_font("Arial", size=12)
                    elif elem_type == "paragraph":
                        pdf.multi_cell(0, 10, text)
                    elif elem_type == "list":
                        for item in element.get("items", []):
                            pdf.cell(0, 10, f"  • {item}", ln=True)
            
            ensure_dir(Path(filepath).parent)
            pdf.output(filepath)
            logger.info(f"Created PDF: {filepath}")
            return filepath
            
        except ImportError:
            raise ImportError("No PDF library available. Install: pip install reportlab fpdf")


# ============================================================================
# EXCEL CREATOR
# ============================================================================

class ExcelCreator:
    """Create Excel spreadsheets (.xlsx)."""
    
    @staticmethod
    def create(
        filepath: str,
        data: Union[Dict, List],
        sheet_name: str = "Sheet1",
    ) -> str:
        """
        Create an Excel file.
        
        Args:
            filepath: Output path
            data: Either:
                  - Dict with sheet names as keys and list of rows as values
                  - List of dicts (each dict is a row)
                  - List of lists (each list is a row)
            sheet_name: Default sheet name (if data is a list)
        
        Returns:
            Path to created file
        """
        # Try openpyxl
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.utils import get_column_letter
            
            wb = Workbook()
            
            # Process data
            if isinstance(data, dict):
                # Multiple sheets
                first_sheet = True
                for sheet_name, sheet_data in data.items():
                    if first_sheet:
                        ws = wb.active
                        ws.title = sheet_name
                        first_sheet = False
                    else:
                        ws = wb.create_sheet(title=sheet_name)
                    
                    ExcelCreator._write_sheet(ws, sheet_data)
            else:
                # Single sheet
                ws = wb.active
                ws.title = sheet_name
                ExcelCreator._write_sheet(ws, data)
            
            ensure_dir(Path(filepath).parent)
            wb.save(filepath)
            logger.info(f"Created Excel file: {filepath}")
            return filepath
            
        except ImportError:
            pass
        
        # Try pandas
        try:
            import pandas as pd
            
            if isinstance(data, dict):
                with pd.ExcelWriter(filepath) as writer:
                    for sheet_name, sheet_data in data.items():
                        df = pd.DataFrame(sheet_data)
                        df.to_excel(writer, sheet_name=sheet_name, index=False)
            else:
                df = pd.DataFrame(data)
                df.to_excel(filepath, sheet_name=sheet_name, index=False)
            
            logger.info(f"Created Excel file: {filepath}")
            return filepath
            
        except ImportError:
            raise ImportError("No Excel library available. Install: pip install openpyxl pandas")
    
    @staticmethod
    def _write_sheet(ws, data):
        """Write data to a worksheet."""
        from openpyxl.styles import Font, PatternFill
        
        if not data:
            return
        
        # Check if list of dicts
        if isinstance(data[0], dict):
            headers = list(data[0].keys())
            
            # Write headers
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
            
            # Write data
            for row_idx, row_data in enumerate(data, 2):
                for col_idx, header in enumerate(headers, 1):
                    ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ""))
        else:
            # List of lists - first row is headers
            for row_idx, row_data in enumerate(data, 1):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    if row_idx == 1:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")


# ============================================================================
# POWERPOINT CREATOR
# ============================================================================

class PowerPointCreator:
    """Create PowerPoint presentations (.pptx)."""
    
    @staticmethod
    def create(
        filepath: str,
        slides: List[Dict],
        title: str = None,
        author: str = "GroundZero AI",
    ) -> str:
        """
        Create a PowerPoint presentation.
        
        Args:
            filepath: Output path
            slides: List of slide definitions:
                    [{"type": "title", "title": "...", "subtitle": "..."},
                     {"type": "content", "title": "...", "content": ["bullet 1", "bullet 2"]},
                     {"type": "two_column", "title": "...", "left": [...], "right": [...]},
                     {"type": "table", "title": "...", "headers": [...], "rows": [...]}]
            title: Presentation title
            author: Author name
        
        Returns:
            Path to created file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation()
        prs.core_properties.author = author
        if title:
            prs.core_properties.title = title
        
        for slide_data in slides:
            slide_type = slide_data.get("type", "content")
            
            if slide_type == "title":
                # Title slide
                layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get("title", "")
                if slide.placeholders[1]:
                    slide.placeholders[1].text = slide_data.get("subtitle", "")
            
            elif slide_type == "content":
                # Bullet points slide
                layout = prs.slide_layouts[1]  # Title and content layout
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get("title", "")
                
                body = slide.placeholders[1]
                tf = body.text_frame
                
                for i, item in enumerate(slide_data.get("content", [])):
                    if i == 0:
                        tf.text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = slide_data.get("level", 0)
            
            elif slide_type == "section":
                # Section header
                layout = prs.slide_layouts[2]  # Section header layout
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get("title", "")
            
            elif slide_type == "blank":
                # Blank slide
                layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(layout)
        
        ensure_dir(Path(filepath).parent)
        prs.save(filepath)
        logger.info(f"Created PowerPoint: {filepath}")
        return filepath


# ============================================================================
# CSV CREATOR
# ============================================================================

class CSVCreator:
    """Create CSV files."""
    
    @staticmethod
    def create(
        filepath: str,
        data: List,
        headers: List[str] = None,
        delimiter: str = ",",
    ) -> str:
        """
        Create a CSV file.
        
        Args:
            filepath: Output path
            data: List of rows (lists or dicts)
            headers: Column headers (auto-detected if data is list of dicts)
            delimiter: Field delimiter
        
        Returns:
            Path to created file
        """
        ensure_dir(Path(filepath).parent)
        
        with open(filepath, 'w', newline='', encoding='utf-8') as f:
            if data and isinstance(data[0], dict):
                # List of dicts
                headers = headers or list(data[0].keys())
                writer = csv.DictWriter(f, fieldnames=headers, delimiter=delimiter)
                writer.writeheader()
                writer.writerows(data)
            else:
                # List of lists
                writer = csv.writer(f, delimiter=delimiter)
                if headers:
                    writer.writerow(headers)
                writer.writerows(data)
        
        logger.info(f"Created CSV: {filepath}")
        return filepath


# ============================================================================
# MARKDOWN CREATOR
# ============================================================================

class MarkdownCreator:
    """Create Markdown files."""
    
    @staticmethod
    def create(
        filepath: str,
        content: Union[str, List[Dict]],
        title: str = None,
    ) -> str:
        """
        Create a Markdown file.
        
        Args:
            filepath: Output path
            content: String or list of elements
            title: Document title
        
        Returns:
            Path to created file
        """
        lines = []
        
        if title:
            lines.append(f"# {title}\n")
        
        if isinstance(content, str):
            lines.append(content)
        else:
            for element in content:
                elem_type = element.get("type", "paragraph")
                
                if elem_type == "heading":
                    level = element.get("level", 1)
                    prefix = "#" * level
                    lines.append(f"{prefix} {element.get('text', '')}\n")
                
                elif elem_type == "paragraph":
                    lines.append(f"{element.get('text', '')}\n")
                
                elif elem_type == "code":
                    lang = element.get("language", "")
                    code = element.get("text", "")
                    lines.append(f"```{lang}\n{code}\n```\n")
                
                elif elem_type == "list":
                    for item in element.get("items", []):
                        lines.append(f"- {item}")
                    lines.append("")
                
                elif elem_type == "numbered_list":
                    for i, item in enumerate(element.get("items", []), 1):
                        lines.append(f"{i}. {item}")
                    lines.append("")
                
                elif elem_type == "table":
                    headers = element.get("headers", [])
                    rows = element.get("rows", [])
                    
                    if headers:
                        lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    
                    for row in rows:
                        lines.append("| " + " | ".join(str(c) for c in row) + " |")
                    lines.append("")
                
                elif elem_type == "quote":
                    for line in element.get("text", "").split("\n"):
                        lines.append(f"> {line}")
                    lines.append("")
                
                elif elem_type == "hr":
                    lines.append("---\n")
        
        ensure_dir(Path(filepath).parent)
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))
        
        logger.info(f"Created Markdown: {filepath}")
        return filepath


# ============================================================================
# UNIFIED FILE CREATOR
# ============================================================================

class FileCreator:
    """
    Unified file creation interface.
    """
    
    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else get_data_path("workspace", "outputs")
        ensure_dir(self.output_dir)
    
    def create(
        self,
        filename: str,
        content: Any,
        file_type: str = None,
        **kwargs
    ) -> str:
        """
        Create a file of any supported type.
        
        Args:
            filename: Output filename
            content: File content
            file_type: File type (auto-detected from extension if None)
            **kwargs: Additional arguments for specific file types
        
        Returns:
            Path to created file
        """
        filepath = self.output_dir / filename
        
        # Auto-detect type from extension
        if file_type is None:
            ext = Path(filename).suffix.lower()
            type_map = {
                '.docx': 'word',
                '.doc': 'word',
                '.pdf': 'pdf',
                '.xlsx': 'excel',
                '.xls': 'excel',
                '.pptx': 'powerpoint',
                '.csv': 'csv',
                '.tsv': 'csv',
                '.md': 'markdown',
                '.txt': 'text',
                '.json': 'json',
                '.html': 'html',
            }
            file_type = type_map.get(ext, 'text')
        
        # Create file
        if file_type == 'word':
            return WordCreator.create(str(filepath), content, **kwargs)
        
        elif file_type == 'pdf':
            return PDFCreator.create(str(filepath), content, **kwargs)
        
        elif file_type == 'excel':
            return ExcelCreator.create(str(filepath), content, **kwargs)
        
        elif file_type == 'powerpoint':
            return PowerPointCreator.create(str(filepath), content, **kwargs)
        
        elif file_type == 'csv':
            if isinstance(content, list):
                return CSVCreator.create(str(filepath), content, **kwargs)
            else:
                raise ValueError("CSV content must be a list")
        
        elif file_type == 'markdown':
            return MarkdownCreator.create(str(filepath), content, **kwargs)
        
        elif file_type == 'json':
            with open(filepath, 'w') as f:
                json.dump(content, f, indent=2, default=str)
            return str(filepath)
        
        else:  # text, html, etc.
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(str(content))
            return str(filepath)
    
    def create_word(self, filename: str, content, **kwargs) -> str:
        """Create Word document."""
        return self.create(filename, content, 'word', **kwargs)
    
    def create_pdf(self, filename: str, content, **kwargs) -> str:
        """Create PDF."""
        return self.create(filename, content, 'pdf', **kwargs)
    
    def create_excel(self, filename: str, data, **kwargs) -> str:
        """Create Excel file."""
        return self.create(filename, data, 'excel', **kwargs)
    
    def create_powerpoint(self, filename: str, slides, **kwargs) -> str:
        """Create PowerPoint."""
        return self.create(filename, slides, 'powerpoint', **kwargs)
    
    def create_csv(self, filename: str, data, **kwargs) -> str:
        """Create CSV."""
        return self.create(filename, data, 'csv', **kwargs)
    
    def create_markdown(self, filename: str, content, **kwargs) -> str:
        """Create Markdown."""
        return self.create(filename, content, 'markdown', **kwargs)


# Export
__all__ = [
    'WordCreator', 'PDFCreator', 'ExcelCreator', 'PowerPointCreator',
    'CSVCreator', 'MarkdownCreator', 'FileCreator',
]
